#!/usr/bin/env python3

import re
from dataclasses import dataclass
from typing import Dict, List, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

@dataclass
class Node:
    id: str
    label: str
    shape_type: str  # 'rectangle', 'diamond', 'oval'
    color: str  # 'blue', 'purple', 'green', 'red'
    x: float = 0
    y: float = 0

@dataclass
class Connection:
    from_node: str
    to_node: str
    label: Optional[str] = None

class MermaidParser:
    def __init__(self, mermaid_content: str):
        self.content = mermaid_content
        self.nodes: Dict[str, Node] = {}
        self.connections: List[Connection] = []
        
    def parse(self):
        lines = self.content.split('\n')
        
        # First pass: extract all node definitions from all lines
        for line in lines:
            line = line.strip()
            if not line or line.startswith('```') or line.startswith('flowchart') or line.startswith('subgraph') or line.startswith('end') or line.startswith('%%') or line.startswith('style'):
                continue
            
            # Extract node definitions even from connection lines
            self._extract_nodes_from_line(line)
        
        # Second pass: extract connections
        for line in lines:
            line = line.strip()
            if not line or line.startswith('```') or line.startswith('flowchart') or line.startswith('subgraph') or line.startswith('end') or line.startswith('%%') or line.startswith('style'):
                continue
                
            if '-->' in line:
                self._parse_connection(line)
    
    def _extract_nodes_from_line(self, line: str):
        # Extract all node definitions from a line, even if it contains connections
        patterns = [
            (r'(\w+)\(\[([^\]]+)\]\)', 'oval'),      # Start([...])
            (r'(\w+)\{([^}]+)\}', 'diamond'),         # RunTests1{...}
            (r'(\w+)\["([^"]+)"\]', 'rectangle'),    # RunLint["..."]
            (r'(\w+)\[([^\]]+)\]', 'rectangle')       # Other rectangles
        ]
        
        for pattern, shape_type in patterns:
            matches = re.findall(pattern, line)
            for match in matches:
                node_id, label = match
                if node_id not in self.nodes:  # Avoid duplicates
                    # Clean up label
                    label = label.replace('<br/>', '\n').replace('\\n', '\n')
                    
                    # Determine color based on label content
                    color = 'blue'  # default
                    if 'CLAUDE' in label:
                        color = 'purple'
                    elif 'DONE' in label or 'ALL DONE' in label:
                        color = 'green'
                    elif 'EXIT' in label or 'ABORT' in label:
                        color = 'red'
                        
                    self.nodes[node_id] = Node(node_id, label, shape_type, color)
    
    def _parse_connection(self, line: str):
        # Parse connections with optional labels
        if '|' in line and '-->' in line:
            # Connection with label: A -->|LABEL| B
            match = re.search(r'(\w+)\s*-->\s*\|([^|]+)\|\s*(\w+)', line)
            if match:
                from_node, label, to_node = match.groups()
                # Only add connection if both nodes exist
                if from_node in self.nodes and to_node in self.nodes:
                    self.connections.append(Connection(from_node, to_node, label.strip()))
        else:
            # Simple connection: A --> B
            match = re.search(r'(\w+)\s*-->\s*(\w+)', line)
            if match:
                from_node, to_node = match.groups()
                # Only add connection if both nodes exist
                if from_node in self.nodes and to_node in self.nodes:
                    self.connections.append(Connection(from_node, to_node))

class FlowchartLayouter:
    def __init__(self, nodes: Dict[str, Node], connections: List[Connection]):
        self.nodes = nodes
        self.connections = connections
        self.levels: Dict[str, int] = {}
        
    def calculate_layout(self, slide_width: float, slide_height: float):
        # Calculate levels (depth from start)
        self._calculate_levels()
        
        # Group nodes by level
        level_groups: Dict[int, List[str]] = {}
        for node_id, level in self.levels.items():
            if level not in level_groups:
                level_groups[level] = []
            level_groups[level].append(node_id)
        
        # Position nodes
        max_level = max(level_groups.keys()) if level_groups else 0
        level_width = slide_width / (max_level + 1) if max_level > 0 else slide_width
        
        for level, node_ids in level_groups.items():
            x = level * level_width + level_width / 2
            node_count = len(node_ids)
            
            for i, node_id in enumerate(node_ids):
                y = slide_height * (i + 1) / (node_count + 1)
                self.nodes[node_id].x = x
                self.nodes[node_id].y = y
    
    def _calculate_levels(self):
        # Find start nodes (nodes with no incoming connections)
        incoming = set()
        for conn in self.connections:
            incoming.add(conn.to_node)
        
        start_nodes = [node_id for node_id in self.nodes.keys() if node_id not in incoming]
        
        # If no start nodes found, use the first node
        if not start_nodes:
            start_nodes = [list(self.nodes.keys())[0]] if self.nodes else []
        
        # BFS to assign levels
        queue = [(node_id, 0) for node_id in start_nodes]
        visited = set()
        
        while queue:
            node_id, level = queue.pop(0)
            if node_id in visited or node_id not in self.nodes:
                continue
                
            visited.add(node_id)
            self.levels[node_id] = level
            
            # Add connected nodes to queue
            for conn in self.connections:
                if conn.from_node == node_id and conn.to_node not in visited and conn.to_node in self.nodes:
                    queue.append((conn.to_node, level + 1))
        
        # Assign level 0 to any unvisited nodes
        for node_id in self.nodes.keys():
            if node_id not in self.levels:
                self.levels[node_id] = 0

class PowerPointGenerator:
    def __init__(self):
        self.prs = Presentation()
        
    def create_flowchart_slide(self, nodes: Dict[str, Node], connections: List[Connection]):
        # Create slide
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Calculate layout
        slide_width = Inches(10)
        slide_height = Inches(7.5)
        layouter = FlowchartLayouter(nodes, connections)
        layouter.calculate_layout(slide_width.inches, slide_height.inches)
        
        # Draw nodes
        shape_map = {}
        for node in nodes.values():
            shape = self._create_node_shape(slide, node)
            if shape:  # Only add if shape creation succeeded
                shape_map[node.id] = shape
        
        # Draw connections
        for conn in connections:
            if (conn.from_node in shape_map and conn.to_node in shape_map and 
                shape_map[conn.from_node] is not None and shape_map[conn.to_node] is not None):
                self._create_connection(slide, shape_map[conn.from_node], shape_map[conn.to_node], conn.label)
        
        return slide
    
    def _create_node_shape(self, slide, node: Node):
        try:
            # Convert position to PowerPoint coordinates
            left = Inches(max(0.1, node.x - 0.75))  # Ensure minimum margin
            top = Inches(max(0.8, node.y - 0.4))    # Leave space for title
            width = Inches(1.5)
            height = Inches(0.8)
            
            # Choose shape type
            if node.shape_type == 'diamond':
                shape_type = MSO_SHAPE.DIAMOND
            elif node.shape_type == 'oval':
                shape_type = MSO_SHAPE.OVAL
            else:
                shape_type = MSO_SHAPE.RECTANGLE
            
            # Create shape
            shape = slide.shapes.add_shape(shape_type, left, top, width, height)
            
            # Set colors
            fill = shape.fill
            fill.solid()
            
            if node.color == 'blue':
                fill.fore_color.rgb = RGBColor(66, 165, 245)  # Blue
            elif node.color == 'purple':
                fill.fore_color.rgb = RGBColor(123, 31, 162)  # Purple
            elif node.color == 'green':
                fill.fore_color.rgb = RGBColor(46, 125, 50)   # Green
            elif node.color == 'red':
                fill.fore_color.rgb = RGBColor(198, 40, 40)   # Red
            
            # Set text
            text_frame = shape.text_frame
            text_frame.text = node.label
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Set text color to white and ensure text fits
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.size = Pt(8)  # Slightly smaller to fit better
            
            # Adjust text frame margins
            text_frame.margin_left = Inches(0.05)
            text_frame.margin_right = Inches(0.05)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            
            return shape
        except Exception as e:
            print(f"Warning: Could not create shape for node {node.id}: {e}")
            return None
    
    def _create_connection(self, slide, from_shape, to_shape, label: Optional[str] = None):
        try:
            # Calculate connection points
            from_x = from_shape.left + from_shape.width / 2
            from_y = from_shape.top + from_shape.height / 2
            to_x = to_shape.left + to_shape.width / 2
            to_y = to_shape.top + to_shape.height / 2
            
            # Create connector
            connector = slide.shapes.add_connector(
                1,  # Straight connector
                from_x, from_y, to_x, to_y
            )
            
            # Style the connector
            line = connector.line
            line.color.rgb = RGBColor(0, 0, 0)
            line.width = Pt(2)
            
            # Add label if provided
            if label:
                mid_x = (from_x + to_x) / 2
                mid_y = (from_y + to_y) / 2
                
                text_box = slide.shapes.add_textbox(
                    mid_x - Inches(0.3), mid_y - Inches(0.1),
                    Inches(0.6), Inches(0.2)
                )
                text_frame = text_box.text_frame
                text_frame.text = label
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Style label text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(8)
                        run.font.color.rgb = RGBColor(0, 0, 0)
        except Exception as e:
            print(f"Warning: Could not create connection: {e}")
    
    def save(self, filename: str):
        self.prs.save(filename)

def main():
    # Read the mermaid diagram
    with open('/Users/llewellynfalco/Github/ExploratoryTestingAI/internal_documentation/linting.diagram.md', 'r') as f:
        content = f.read()
    
    # Extract mermaid content
    mermaid_match = re.search(r'```mermaid\n(.*?)\n```', content, re.DOTALL)
    if not mermaid_match:
        print("No mermaid diagram found in the file!")
        return
    
    mermaid_content = mermaid_match.group(1)
    
    # Parse the diagram
    parser = MermaidParser(mermaid_content)
    parser.parse()
    
    print(f"Parsed {len(parser.nodes)} nodes and {len(parser.connections)} connections")
    
    # Generate PowerPoint
    generator = PowerPointGenerator()
    slide = generator.create_flowchart_slide(parser.nodes, parser.connections)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Linting Automation Process Flow"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Save the presentation
    output_file = '/Users/llewellynfalco/Github/ExploratoryTestingAI/linting_flowchart.pptx'
    generator.save(output_file)
    print(f"PowerPoint presentation saved as: {output_file}")

if __name__ == "__main__":
    main()
